#!/usr/bin/env python3
"""Costruisce la presentazione 'WC muon veto — rate estimate' (inglese).
Narrativa: prototipo testbeam (eff. generale + eff vs L) -> veto reale ->
stima rate (mu + gamma) -> riduzione della componente gamma.
Immagini prese dallo scratchpad (ppt/). Rilanciabile."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import os

PNG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "pptx_assets")
OUT = "/Users/benussi/Testbeam2026_WC_unified/WC_muon_veto_rate_estimate.pptx"
BLU=RGBColor(0x1F,0x3A,0x5F); GRIG=RGBColor(0x44,0x44,0x44); ROSSO=RGBColor(0xB2,0x22,0x22); VERDE=RGBColor(0x1a,0x7a,0x3a)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def title(s,text,sub=None):
    tb=s.shapes.add_textbox(Inches(0.45),Inches(0.2),Inches(12.4),Inches(1.0))
    p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text=text
    r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=BLU
    if sub:
        p2=tb.text_frame.add_paragraph(); r2=p2.add_run(); r2.text=sub
        r2.font.size=Pt(15); r2.font.color.rgb=GRIG

def bullets(s,items,top=1.4,left=0.6,width=12.2,size=18):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.6)); tf=tb.text_frame; tf.word_wrap=True
    for i,(t,lvl,bold,col) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.level=lvl
        r=p.add_run(); r.text=("• " if lvl==0 else "– ")+t
        r.font.size=Pt(size-2*lvl); r.font.bold=bold; r.font.color.rgb=col or GRIG; p.space_after=Pt(7)

def img(s,name,top=1.4,maxh=5.5,maxw=12.4,left=None):
    path=os.path.join(PNG,name); im=Image.open(path); ar=im.size[0]/im.size[1]
    W=maxw; H=W/ar
    if H>maxh: H=maxh; W=H*ar
    L=(13.333-W)/2 if left is None else left
    s.shapes.add_picture(path,Inches(L),Inches(top),Inches(W),Inches(H))

def note(s,text,top=6.98):
    tb=s.shapes.add_textbox(Inches(0.45),Inches(top),Inches(12.4),Inches(0.4))
    r=tb.text_frame.paragraphs[0].add_run(); r.text=text
    r.font.size=Pt(12); r.font.italic=True; r.font.color.rgb=GRIG

def table(s,rows,left=1.2,top=1.7,width=10.9,colw=None,fs=15,head=True):
    nr=len(rows); nc=len(rows[0])
    gt=s.shapes.add_table(nr,nc,Inches(left),Inches(top),Inches(width),Inches(0.4*nr)).table
    if colw:
        for j,w in enumerate(colw): gt.columns[j].width=Inches(w)
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.text=str(val)
            pr=c.text_frame.paragraphs[0]; pr.runs[0].font.size=Pt(fs)
            pr.runs[0].font.bold=(i==0 and head)
            if i==0 and head: pr.runs[0].font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
    return gt

B=lambda t,l=0,b=False,c=None:(t,l,b,c)

# ---------- 1 TITLE ----------
s=prs.slides.add_slide(BLANK)
tb=s.shapes.add_textbox(Inches(0.8),Inches(2.1),Inches(11.7),Inches(3.0))
p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text="Water-Cherenkov Muon Veto for CUPID"
r.font.size=Pt(40); r.font.bold=True; r.font.color.rgb=BLU
p2=tb.text_frame.add_paragraph(); r2=p2.add_run()
r2.text="From the testbeam prototype to the rate estimate on the full detector"
r2.font.size=Pt(21); r2.font.color.rgb=GRIG
p3=tb.text_frame.add_paragraph(); r3=p3.add_run()
r3.text="Measured efficiency  ->  muon/gamma rate  ->  how we suppress the gamma component"
r3.font.size=Pt(16); r3.font.color.rgb=VERDE
p4=tb.text_frame.add_paragraph(); r4=p4.add_run(); r4.text="CUPID muon veto — LNGS"
r4.font.size=Pt(15); r4.font.color.rgb=GRIG

# ---------- 2 MOTIVATION ----------
s=prs.slides.add_slide(BLANK); title(s,"Goal & requirements")
bullets(s,[
 B("A Water-Cherenkov (WC) detector will act as the muon veto of CUPID at LNGS",0,True,BLU),
 B("Cosmic muons induce backgrounds that can mimic the 0nbb signal in the bolometers",1),
 B("Muons must be tagged with high efficiency and flagged in a veto time window",1),
 B("Two competing requirements",0,True,BLU),
 B("High muon efficiency  ->  tag as many muons as possible",1,False,VERDE),
 B("Low gamma fake rate  ->  environmental gammas must NOT fire the veto (they cause dead time)",1,False,ROSSO),
 B("This talk: measured prototype performance, scaled to the full veto, and the resulting rate",0,True,BLU),
])

# ---------- 3 PROTOTYPE ----------
s=prs.slides.add_slide(BLANK); title(s,"The prototype module (tested at the testbeam)")
bullets(s,[
 B("Single WC module: HDPE tank 98 x 36 x 12 cm,  10 cm water depth",0,True,BLU),
 B("Inner Teflon (PTFE) reflector (0.91-0.97)",1),
 B("37 WLS fibers (BCF-9995XL) collect the Cherenkov light",1),
 B("Fibers routed to a bundle read by 2 PMTs (A, B)",1),
 B("Muon tag = A OR B above threshold",0,True,BLU),
 B("Charge plane (A,B) used for muon/gamma separation",1),
 B("Beam: 500 MeV e / mu; positions & angles scanned",1),
],top=1.5,left=0.55,width=5.3,size=16)
img(s,"tb_module_schematic.png",top=1.7,maxh=4.6,maxw=7.0,left=6.0)
note(s,"Top view (fiber layout, PMT A/B) and side view (straight-in-water + Bezier routing to the vertical bundle).")

# ---------- 4 GENERAL EFFICIENCY ----------
s=prs.slides.add_slide(BLANK)
title(s,"Testbeam result: muon efficiency ~ 98-99%","2D map across the module face (theta=0 deg): the two-PMT tag (A OR B) is uniform and near-full.")
img(s,"tb_eff_2Dmap.png",top=1.75,maxh=4.9)
note(s,"Prototype data, module 0. Total efficiency 92-100% over the whole active surface (mostly 98-99%).")

# ---------- 5 EFF vs L ----------
s=prs.slides.add_slide(BLANK)
title(s,"Testbeam result: efficiency vs track length","DATA confirm ~99% at all measured lengths; MC reproduces the data and predicts the roll-off only at very short tracks.")
img(s,"tb_eff_vs_L.png",top=1.75,maxh=4.9)
note(s,"DATA (black) ~99% for L >= 5 cm. MC (red) matches and shows the drop only for L < ~4 cm (grazing muons, little light).")

# ---------- 6 FROM PROTOTYPE TO FULL VETO ----------
s=prs.slides.add_slide(BLANK)
title(s,"From the prototype to the full CUPID muon veto","Real engineering CAD (20260713-Assy-WC-Y): box 3.78 x 3.87 x 1.64 m, 48 water tanks, 96 PMTs.")
img(s,"real_geometry.png",top=1.5,maxh=5.1,maxw=8.3,left=0.25)
bullets(s,[
 B("Real CAD, real holes",0,True,BLU),
 B("48 tanks: 41 instrumented + 7 corner (non-sensitive)",1),
 B("96 PMTs = 2 per tank (black markers)",1),
 B("Circular hole in the bottom + wall ports",1),
 B("Water = tank inner cavity (100 mm)",1),
 B("84 active sectors (<=40 cm) in the veto model",1),
 B("Prototype efficiency vs L assigned per sector",1),
],top=1.7,left=8.8,width=4.3,size=14)
note(s,"Container mesh with the real holes + PMT positions from the CAD; water is the inner cavity of each tank.")

# ---------- 7 MUON EFFICIENCY ON FULL DETECTOR ----------
s=prs.slides.add_slide(BLANK)
title(s,"Muon coverage of the full detector","Ray-tracing cosmic muons through the 84 real water sectors, prototype light response sampled per sector.")
bullets(s,[
 B("Every muon crosses >= 2 water sectors (median 2) -> multiple chances to tag",0,True,BLU),
 B("Box tagging efficiency  eps_box (A OR B) ~ 99%   (WBC2 98.8% / WBC3 99.0%)",0,True,VERDE),
 B("The ~1% lost are corner-skimmers with very short path (little light everywhere)",1),
 B("Muon rate on the box:  ~ 4.5 mHz  =  ~ 390 muons/day",0,True,BLU),
 B("from Mei & Hime flux (3e-8 cm-2 s-1) x shadow area ~15 m2",1),
 B("The real geometry (real holes, water cavities) CONFIRMS the estimate: eps_box stays ~99%",0,True,VERDE),
])

# ---------- 8 GAMMA COMPONENT ----------
s=prs.slides.add_slide(BLANK)
title(s,"The gamma component","Environmental gammas hit the large surface at high rate and can fake a muon.")
img(s,"gamma_vs_mu_charge.png",top=1.65,maxh=4.0)
bullets(s,[
 B("Gamma flux Hall A (measured, arXiv:1101.5298): 0.35 cm-2 s-1  ->  ~ 190 kHz on the ~54 m2 box surface",0,True,ROSSO),
 B("Gammas make little light (low charge) -> separable from muons in the two-PMT charge plane",0,True,VERDE),
],top=5.35,size=15)

# ---------- 9 REDUCING GAMMA ----------
s=prs.slides.add_slide(BLANK)
title(s,"Suppressing the gamma component: muon/gamma discrimination","Cut in the two-PMT charge plane (A,B): keep large-charge (muon), reject the low-charge gamma corner.")
img(s,"cuts_illustration.png",top=1.6,maxh=4.7)
bullets(s,[
 B("A+B (sum) and radial cut sqrt(A^2+B^2) are equivalent and hardware-simple; both beat A AND B coincidence",0,True,VERDE),
 B("Threshold tuned to a target gamma fake rate (e.g. 0.5 Hz) while keeping muon eps ~ 99%",0,True,BLU),
],top=6.15,size=14)

# ---------- 10 RATE ESTIMATE ----------
s=prs.slides.add_slide(BLANK)
title(s,"Estimated rate on the muon veto","Total veto trigger rate = tagged muons + residual gamma fakes.")
table(s,[
 ["Component","Rate on the box","Note"],
 ["Muons crossing","~ 4.5 mHz  (390/day)","Mei & Hime x shadow area"],
 ["Muons tagged (R_mu,tag)","~ 4.5 mHz","= R_mu x eps_box (~99%)"],
 ["Gammas raw","~ 190 kHz","0.35 cm-2 s-1 x 54 m2"],
 ["Gamma FAKE (after cut)","0.5 Hz  (tunable)","tuned via the A+B / radial threshold"],
 ["TOTAL veto trigger","~ 0.5 Hz","gamma-dominated"],
],top=1.9,fs=15,colw=[3.3,3.4,4.2])
bullets(s,[
 B("The veto trigger rate is set by the gamma fakes, NOT by the muons (muons are 4 orders of magnitude rarer)",0,True,ROSSO),
 B("Lowering the fake target (tighter cut) directly lowers the veto rate — at fixed muon efficiency",0,True,VERDE),
],top=4.9,size=15)

# ---------- 11 DEAD TIME PREVIEW ----------
s=prs.slides.add_slide(BLANK)
title(s,"Dead time — preview","dead time = (gamma fake + true muons) x tau_veto.  tau_veto to be defined with trigger/DAQ.")
table(s,[
 ["tau_veto","1 ms","5 ms","10 ms","20 ms","50 ms"],
 ["dead time (gamma fake 0.5 Hz)","0.05%","0.25%","0.51%","1.0%","2.5%"],
 ["good events lost / yr (LD 5 Hz)","80 k","0.4 M","0.8 M","1.6 M","4.0 M"],
],top=1.9,fs=14,colw=[4.2,1.7,1.7,1.7,1.7,1.7])
bullets(s,[
 B("True-muon dead time is negligible (~5e-4 % at 1 ms): the muon veto is essentially 'free'",0,True,VERDE),
 B("All the dead time comes from gamma fakes -> for long gates (tens of ms) the gamma rejection is the critical knob",0,True,ROSSO),
 B("e.g. at tau=50 ms, keeping dead time < 0.1% needs the gamma fake rate ~ 0.02 Hz (a 25x tighter cut)",1),
],top=4.4,size=15)
note(s,"tau_veto is a trigger/DAQ choice; the table shows the trend, not a final number.")

# ---------- 12 SUMMARY ----------
s=prs.slides.add_slide(BLANK); title(s,"Summary & outlook")
bullets(s,[
 B("Prototype demonstrated ~99% muon efficiency at the testbeam (two-PMT A OR B, nominal crossing)",0,True,VERDE),
 B("Efficiency is high and flat vs track length; drops only for very short grazing tracks",1),
 B("Validated on the real CAD (84 water sectors, real holes): eps_box ~ 99%",0,True,VERDE),
 B("Muon rate ~ 390/day; the muon veto adds negligible dead time by itself",1),
 B("The rate on the veto is dominated by environmental gammas",0,True,ROSSO),
 B("Two-PMT charge-plane cut (A+B / radial) suppresses gammas to a tunable fake rate (~0.5 Hz)",1),
 B("Next steps",0,True,BLU),
 B("Fix tau_veto with trigger/DAQ; optimize the gamma-fake vs dead-time working point",1),
 B("Finalize the photosensor (H10721-210) and the fiber routing (WBC2/WBC3)",1),
])

prs.save(OUT)
print("salvato:",OUT,"—",len(prs.slides._sldIdLst),"slide")
