from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import os

PNG="/Users/benussi/Testbeam2026_WC_unified/Analysis_script_background/output_final"
OUT="/Users/benussi/Testbeam2026_WC_unified/WC_muon_veto_CUPID_overview.pptx"
BLU=RGBColor(0x1F,0x3A,0x5F); GRIG=RGBColor(0x44,0x44,0x44); ROSSO=RGBColor(0xB2,0x22,0x22); VERDE=RGBColor(0x1a,0x7a,0x3a)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def title(s, text, sub=None):
    tb=s.shapes.add_textbox(Inches(0.45),Inches(0.2),Inches(12.4),Inches(1.0))
    p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text=text
    r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=BLU
    if sub:
        p2=tb.text_frame.add_paragraph(); r2=p2.add_run(); r2.text=sub
        r2.font.size=Pt(15); r2.font.color.rgb=GRIG

def bullets(s, items, top=1.35, left=0.6, width=12.2, size=18):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.6)); tf=tb.text_frame; tf.word_wrap=True
    for i,(t,lvl,bold,col) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.level=lvl
        r=p.add_run(); r.text=("• " if lvl==0 else "– ")+t
        r.font.size=Pt(size-2*lvl); r.font.bold=bold; r.font.color.rgb=col or GRIG; p.space_after=Pt(7)

def img(s, name, top=1.4, maxh=5.6, maxw=12.4):
    path=os.path.join(PNG,name); im=Image.open(path); ar=im.size[0]/im.size[1]
    W=maxw; H=W/ar
    if H>maxh: H=maxh; W=H*ar
    s.shapes.add_picture(path, Inches((13.333-W)/2), Inches(top), Inches(W), Inches(H))

def note(s, text, top=6.95):
    tb=s.shapes.add_textbox(Inches(0.45),Inches(top),Inches(12.4),Inches(0.4))
    r=tb.text_frame.paragraphs[0].add_run(); r.text=text
    r.font.size=Pt(12); r.font.italic=True; r.font.color.rgb=GRIG

B=lambda t,l=0,b=False,c=None:(t,l,b,c)

# 1 title
s=prs.slides.add_slide(BLANK)
tb=s.shapes.add_textbox(Inches(0.8),Inches(2.3),Inches(11.7),Inches(2.6))
p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text="A Water-Cherenkov Muon Veto for CUPID"
r.font.size=Pt(40); r.font.bold=True; r.font.color.rgb=BLU
p2=tb.text_frame.add_paragraph(); r2=p2.add_run()
r2.text="Monte-Carlo design study: detector, photosensor, and muon/gamma discrimination"
r2.font.size=Pt(20); r2.font.color.rgb=GRIG
p3=tb.text_frame.add_paragraph(); r3=p3.add_run(); r3.text="Geant4 simulation — LNGS"
r3.font.size=Pt(16); r3.font.color.rgb=GRIG

# 2 context
s=prs.slides.add_slide(BLANK); title(s,"The goal")
bullets(s,[
 B("A Water-Cherenkov (WC) detector will act as the muon veto of the CUPID experiment at LNGS",0,True,BLU),
 B("CUPID searches for neutrinoless double-beta decay with cryogenic bolometers",1),
 B("Cosmic muons induce backgrounds that mimic the signal — they must be tagged and removed",1),
 B("Two competing requirements",0,True,BLU),
 B("Tag as many muons as possible (high signal efficiency)",1),
 B("Do NOT fire on the ambient gamma background — a false tag throws away good physics data (dead time)",1),
 B("This study, entirely in Monte-Carlo, optimises the detector to meet both",0,True,BLU),
], size=19)

# 3 what we studied
s=prs.slides.add_slide(BLANK); title(s,"What we studied — three questions")
bullets(s,[
 B("1.  Detector geometry",0,True,BLU),
 B("How to route the wavelength-shifting fibres to collect the most light with realistic bends",1),
 B("2.  Photosensor",0,True,BLU),
 B("The previous PMT is discontinued — we compared candidates (new PMT vs SiPM)",1),
 B("3.  Selection logic",0,True,BLU),
 B("How to decide, from the two PMT signals, whether an event is a muon or a gamma",1),
 B("All three assessed on the same simulated events, then projected to the full-size experiment",0,False,GRIG),
], size=19)

# 4 geometry
s=prs.slides.add_slide(BLANK); title(s,"Detector: fibre light collection")
bullets(s,[
 B("The light produced by a muon is guided to the PMTs by wavelength-shifting fibres",0),
 B("Optimising the fibre bends increased the collected light by ~25% over the original design",0,True,VERDE),
 B("The chosen bend shape matches the natural shape a real fibre takes when bent — realistic and easy to build",1),
 B("A higher-light variant with more fibres was also tested (small extra gain, kept as a margin option)",1),
 B("Bottom line: the detector is efficient at turning a muon into a measurable signal",0,True,BLU),
], size=19)

# 5 photosensor
s=prs.slides.add_slide(BLANK); title(s,"Photosensor: choice made")
bullets(s,[
 B("Compared a new photomultiplier (H10721-210) against a silicon photomultiplier (SiPM)",0),
 B("Chosen: the new PMT",0,True,VERDE),
 B("Clean, large signals; negligible internal noise",1),
 B("The SiPM collects more light but its internal dark noise lets ~15x more gamma background through",1),
 B("Kept the SiPM option documented for possible future comparison",1),
 B("Bottom line: the new PMT gives the best signal-to-background for a veto",0,True,BLU),
], size=19)

# 6 key plot
s=prs.slides.add_slide(BLANK); title(s,"The key idea: muons vs gammas",
   "Each event gives two signals (PMT A and PMT B). Muons and gammas land in different places.")
img(s,"muon_gamma_cuts_illustration.png", top=1.35, maxh=5.5)
note(s,"Muons (red) light both PMTs → spread across the plane. Gammas (blue) make little light → cluster near the origin.")

# 7 cuts
s=prs.slides.add_slide(BLANK); title(s,"A simple, robust selection")
bullets(s,[
 B("Gammas make little light and sit in the bottom-left corner; muons fill the rest of the plane",0),
 B("We reject the gamma corner with a single threshold on the total light:",0,True,BLU),
 B("‘sum’ cut:  QA + QB above a value   (a straight line)",1),
 B("‘circular’ cut:  QA² + QB² above a value   (a quarter circle)",1),
 B("The two are practically equivalent — the SUM is trivial to build in hardware (analog sum of two PMTs)",0,True,VERDE),
 B("This beats the naive ‘both PMTs must fire’ logic by ~5% in muon efficiency at the same background",0),
], size=18)

# 8 performance
s=prs.slides.add_slide(BLANK); title(s,"Performance at the full experiment",
   "Full detector = a 3.5 x 3.5 x 1.5 m box of sensitive modules (one face has a 1.5 m opening)")
bullets(s,[
 B("At a chosen residual gamma false-tag rate of ~0.5 Hz over the whole detector:",0,True,BLU),
 B("Muon tagging efficiency  ≈ 99%   (a muon crosses two sensitive walls → two chances to be caught)",1,True,VERDE),
 B("Data lost to false gamma tags:  ~0.01%  of exposure (at the expected veto window)",1,True,VERDE),
 B("The residual gamma leakage scales with the true ambient gamma flux — the one input still to be measured on site",0,False,GRIG),
 B("Comparable existing systems (CROSS) accept ~18% dead time — here it is orders of magnitude smaller",0),
], size=18)

# 9 decision map
s=prs.slides.add_slide(BLANK); title(s,"Operating-point map",
   "Good events lost per year vs assumed detector rate and veto window (★ = expected CUPID working point)")
img(s,"decision_2D_map.png", top=1.4, maxh=5.4)

# 10 conclusions
s=prs.slides.add_slide(BLANK); title(s,"Summary")
bullets(s,[
 B("A WC muon veto for CUPID has been designed and optimised in Monte-Carlo",0,True,BLU),
 B("Fibre geometry: +25% light with a realistic, buildable bend",1),
 B("Photosensor: new PMT (H10721-210) — best signal-to-background",1),
 B("Selection: a simple total-light cut (sum of the two PMTs) cleanly separates muons from gammas",1),
 B("Projected performance: ~99% muon efficiency with negligible dead time",1,True,VERDE),
 B("Next steps",0,True,BLU),
 B("Measure the ambient gamma flux at the veto location (the last external input)",1),
 B("A new observable — how many detector sectors light up — could recover the hardest (grazing) muons",1),
 B("Prototype validation at the 2026 test beam",1),
], size=17)

prs.save(OUT)
print("Salvata:", OUT, "-", len(prs.slides._sldIdLst), "slide")
